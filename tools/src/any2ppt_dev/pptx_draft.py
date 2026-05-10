"""Experimental pptx-native assembler. Supports two archetypes (cover, thesis).

This module is intentionally minimal. The Week 4 goal is to discover the real
cost of pptx-native assembly, not to replace image-first generation.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


NEAR_BLACK = RGBColor(0x11, 0x12, 0x13)
SOFT_GRAY = RGBColor(0x55, 0x55, 0x55)
BG_OFFWHITE = RGBColor(0xF4, 0xF1, 0xEA)
BRAND_BLUE = RGBColor(0x25, 0x63, 0xEB)


@dataclass
class Slide:
    slide_id: str
    title: str = ""
    primary_job: str = ""
    core_claim: str = ""
    archetype: str = ""
    presenter: str = ""
    support: list[str] = field(default_factory=list)


SLIDE_HEADING_RE = re.compile(r"^##\s+(\d{2}_[a-z0-9_]+)\s*$", re.MULTILINE)


def parse_storyboard(path: Path) -> list[Slide]:
    text = path.read_text(encoding="utf-8")
    matches = list(SLIDE_HEADING_RE.finditer(text))
    if not matches:
        raise ValueError(f"no slide headings (## NN_slug) found in {path}")
    slides: list[Slide] = []
    for idx, match in enumerate(matches):
        slide_id = match.group(1)
        body_start = match.end()
        body_end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        body = text[body_start:body_end]
        slides.append(_parse_slide_body(slide_id, body))
    return slides


def _parse_slide_body(slide_id: str, body: str) -> Slide:
    slide = Slide(slide_id=slide_id)
    slide.title = _extract_single_field(body, "Title")
    slide.primary_job = _extract_single_field(body, "Primary job")
    slide.core_claim = _extract_single_field(body, "Core claim")
    slide.archetype = _extract_single_field(body, "Archetype")
    slide.presenter = _extract_single_field(body, "Presenter intent")
    slide.support = _extract_support_points(body)
    return slide


def _extract_single_field(body: str, name: str) -> str:
    pattern = rf"^-\s+\*\*{re.escape(name)}\*\*:\s*(.+)$"
    m = re.search(pattern, body, re.MULTILINE | re.IGNORECASE)
    return m.group(1).strip() if m else ""


def _extract_support_points(body: str) -> list[str]:
    pattern = r"^-\s+\*\*Support\s+points?\*\*:\s*\n((?:[ \t]+-\s+[^\n]+\n?)+)"
    m = re.search(pattern, body, re.IGNORECASE | re.MULTILINE)
    if not m:
        return []
    items: list[str] = []
    for line in m.group(1).splitlines():
        stripped = line.strip()
        if not stripped.startswith("-"):
            continue
        items.append(stripped.lstrip("-").strip())
    return items


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size=Pt(24),
    bold: bool = False,
    color=NEAR_BLACK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def render_cover(prs: Presentation, slide_data: Slide) -> None:
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_background(sl, BG_OFFWHITE)

    _add_textbox(
        sl, Inches(0.6), Inches(0.5), Inches(3), Inches(0.5),
        "Any2PPT", font_size=Pt(18), bold=True, color=BRAND_BLUE,
    )

    _add_textbox(
        sl, Inches(0.6), Inches(2.5), Inches(12), Inches(2),
        slide_data.title, font_size=Pt(44), bold=True, color=NEAR_BLACK,
    )

    subtitle = _pick_support(slide_data.support, ("subtitle", "subline"))
    if subtitle:
        _add_textbox(
            sl, Inches(0.6), Inches(4.8), Inches(12), Inches(0.6),
            subtitle, font_size=Pt(20), color=SOFT_GRAY,
        )

    source_line = _pick_support(slide_data.support, ("source",))
    if source_line:
        _add_textbox(
            sl, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
            source_line, font_size=Pt(14), color=SOFT_GRAY,
            align=PP_ALIGN.RIGHT,
        )


def render_thesis(prs: Presentation, slide_data: Slide) -> None:
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_background(sl, BG_OFFWHITE)

    _add_textbox(
        sl, Inches(0.6), Inches(0.4), Inches(12), Inches(1),
        slide_data.title, font_size=Pt(36), bold=True, color=NEAR_BLACK,
    )

    if slide_data.core_claim:
        _add_textbox(
            sl, Inches(0.6), Inches(1.6), Inches(12), Inches(1.4),
            slide_data.core_claim, font_size=Pt(22), color=SOFT_GRAY,
        )

    pillars = slide_data.support[:4]
    n = len(pillars)
    if n == 0:
        return
    block_top = Inches(4.2)
    block_height = Inches(2.0)
    canvas_left = Inches(0.6)
    canvas_width = Inches(12.1)
    gap = Inches(0.3)
    block_width = Emu(int((canvas_width.emu - gap.emu * (n - 1)) / n))
    for i, label in enumerate(pillars):
        left = Emu(canvas_left.emu + i * (block_width.emu + gap.emu))
        shape = sl.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, block_top, block_width, block_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_OFFWHITE
        shape.line.color.rgb = NEAR_BLACK
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = NEAR_BLACK


def render_skipped(prs: Presentation, slide_data: Slide) -> None:
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_background(sl, BG_OFFWHITE)
    _add_textbox(
        sl, Inches(0.6), Inches(0.4), Inches(12), Inches(1),
        slide_data.title or slide_data.slide_id,
        font_size=Pt(32), bold=True, color=NEAR_BLACK,
    )
    note = (
        f"Archetype '{slide_data.archetype}' is not yet supported by the "
        f"pptx-native experimental assembler.\n\n"
        f"Claim: {slide_data.core_claim or '(none)'}"
    )
    _add_textbox(
        sl, Inches(0.6), Inches(2.0), Inches(12), Inches(4),
        note, font_size=Pt(18), color=SOFT_GRAY,
    )


def _pick_support(items: list[str], keywords) -> str:
    for item in items:
        lowered = item.lower()
        for kw in keywords:
            if kw in lowered:
                _, _, value = item.partition(":")
                return value.strip() if value else item
    return ""


def draft(storyboard_path: Path, out_path: Path) -> int:
    storyboard_path = storyboard_path.resolve()
    if not storyboard_path.is_file():
        raise FileNotFoundError(f"storyboard not found: {storyboard_path}")

    slides = parse_storyboard(storyboard_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    counts = {"cover": 0, "thesis": 0, "skipped": 0}
    for s in slides:
        archetype_lower = s.archetype.lower()
        if archetype_lower.startswith("cover"):
            render_cover(prs, s)
            counts["cover"] += 1
        elif archetype_lower.startswith("thesis"):
            render_thesis(prs, s)
            counts["thesis"] += 1
        else:
            render_skipped(prs, s)
            counts["skipped"] += 1

    out_path = out_path.resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote: {out_path}")
    print(
        f"slides: total {len(slides)} "
        f"(cover {counts['cover']}, thesis {counts['thesis']}, skipped {counts['skipped']})"
    )
    return 0
