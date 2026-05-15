from __future__ import annotations

import argparse
import json
import re
import shutil
import zipfile
from pathlib import Path

import yaml


REPO_ROOT = Path(__file__).resolve().parents[3]
DEFAULT_PLUGIN = REPO_ROOT / "plugins" / "deckit"
DEFAULT_RUNS_DIR = REPO_ROOT / "local-runs"
DEFAULT_MARKETPLACE = REPO_ROOT / ".agents" / "plugins" / "marketplace.json"
TEXT_SOURCE_SUFFIXES = {".md", ".markdown", ".txt"}
PRODUCTION_MODES = ("image-first",)


def _load_skill_frontmatter(skill_md: Path) -> dict[str, object]:
    text = skill_md.read_text(encoding="utf-8")
    if not text.startswith("---\n"):
        raise ValueError(f"{skill_md} is missing YAML frontmatter")
    try:
        _, frontmatter, _ = text.split("---\n", 2)
    except ValueError as exc:
        raise ValueError(f"{skill_md} has malformed YAML frontmatter") from exc
    data = yaml.safe_load(frontmatter)
    if not isinstance(data, dict):
        raise ValueError(f"{skill_md} frontmatter is not a mapping")
    return data


def inspect_plugin(plugin: Path) -> int:
    manifest_path = plugin / ".codex-plugin" / "plugin.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    skills_dir = plugin / "skills"
    skill_files = sorted(skills_dir.glob("*/SKILL.md"))

    print(f"plugin: {manifest.get('name')}")
    print(f"version: {manifest.get('version')}")
    print(f"skills: {len(skill_files)}")

    for skill_md in skill_files:
        data = _load_skill_frontmatter(skill_md)
        name = data.get("name")
        description = data.get("description")
        if not name or not description:
            raise ValueError(f"{skill_md} must include name and description")
        print(f"- {name}: {skill_md.parent.relative_to(plugin)}")

    return 0


def inspect_marketplace(marketplace: Path) -> int:
    marketplace = marketplace.resolve()
    data = json.loads(marketplace.read_text(encoding="utf-8"))
    plugins = data.get("plugins")
    if not isinstance(plugins, list):
        raise ValueError(f"{marketplace} must contain a plugins list")

    print(f"marketplace: {data.get('name')}")
    print(f"plugins: {len(plugins)}")

    root = marketplace.parents[2]
    for plugin in plugins:
        name = plugin.get("name")
        source = plugin.get("source", {})
        path_value = source.get("path")
        if not name or not path_value:
            raise ValueError(f"marketplace plugin entry is missing name or source.path: {plugin}")
        plugin_path = (root / path_value).resolve()
        manifest_path = plugin_path / ".codex-plugin" / "plugin.json"
        if not manifest_path.is_file():
            raise FileNotFoundError(f"plugin manifest not found for {name}: {manifest_path}")
        try:
            display_path = plugin_path.relative_to(root)
        except ValueError:
            display_path = plugin_path
        print(f"- {name}: {display_path}")

    return 0


def _slugify(value: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
    slug = re.sub(r"-{2,}", "-", slug)
    if not slug:
        raise ValueError("run name must contain at least one letter or digit")
    return slug


def _ensure_inside(parent: Path, child: Path) -> None:
    parent = parent.resolve()
    child = child.resolve()
    if child != parent and parent not in child.parents:
        raise ValueError(f"path is outside expected directory: {child}")


def _classify_source(source: str) -> str:
    """Return one of: 'url', 'pdf', 'text'."""
    from deckit_dev.ingest import is_url

    if is_url(source):
        return "url"
    suffix = Path(source).suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix in TEXT_SOURCE_SUFFIXES:
        return "text"
    raise ValueError(
        f"unsupported source: {source}. Supported: text files ({sorted(TEXT_SOURCE_SUFFIXES)}), .pdf files, and http(s) URLs."
    )


def _slug_from_source(source: str, kind: str) -> str:
    if kind == "url":
        from urllib.parse import urlparse

        parsed = urlparse(source)
        host = parsed.netloc.replace(".", "-")
        path = parsed.path.strip("/").replace("/", "-") or "index"
        return _slugify(f"{host}-{path}")
    return _slugify(Path(source).stem)


def new_run(
    source: str,
    name: str | None,
    runs_dir: Path,
    force: bool,
    mode: str | None,
    budget: str | None,
) -> int:
    mode = mode or "image-first"
    kind = _classify_source(source)

    if kind == "text":
        source_path = Path(source).resolve()
        if not source_path.is_file():
            raise FileNotFoundError(f"source file does not exist: {source_path}")
    elif kind == "pdf":
        source_path = Path(source).resolve()
        if not source_path.is_file():
            raise FileNotFoundError(f"pdf source does not exist: {source_path}")
    else:
        source_path = None

    run_name = _slugify(name) if name else _slug_from_source(source, kind)
    runs_dir = runs_dir.resolve()
    run_dir = runs_dir / run_name
    _ensure_inside(runs_dir, run_dir)

    if run_dir.exists() and not force:
        raise FileExistsError(f"run already exists: {run_dir}. Use --force to reuse it.")

    artifacts: dict[str, str] = {
        "deck_brief": "work/deck-brief.md",
        "storyboard": "work/storyboard.md",
        "prompt_readme": "prompts/README.md",
        "prompt_files": "prompts/<slide-id>.md",
        "generated_slides": "assets/generated-slides/<slide-id>.png",
    }
    artifacts["dist"] = "dist/"
    artifacts["review"] = "dist/review.md"

    children: list[str] = ["source", "work", "prompts", "assets/generated-slides", "dist"]
    for child in children:
        (run_dir / child).mkdir(parents=True, exist_ok=True)

    if kind == "text":
        target_source = run_dir / "source" / f"input{source_path.suffix.lower()}"
        if target_source.exists() and not force:
            raise FileExistsError(f"source already exists: {target_source}. Use --force to replace it.")
        shutil.copy2(source_path, target_source)
        original_path = str(source_path)
    else:
        target_source = run_dir / "source" / "input.md"
        if target_source.exists() and not force:
            raise FileExistsError(f"source already exists: {target_source}. Use --force to replace it.")
        from deckit_dev.ingest import ingest_pdf, ingest_url

        if kind == "pdf":
            ingest_pdf(source_path, target_source)
            original_path = str(source_path)
        else:
            ingest_url(source, target_source)
            original_path = source

    manifest: dict[str, object] = {
        "name": run_name,
        "source": {
            "original_path": original_path,
            "local_path": str(target_source.relative_to(run_dir)),
            "type": kind,
        },
        "production_mode": mode,
        "budget_mode": budget,
        "artifacts": artifacts,
    }
    (run_dir / "run.json").write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")

    print(f"created run: {run_dir}")
    print(f"source ({kind}): {target_source}")
    print(f"production_mode: {mode}")
    if budget is None:
        print("budget_mode: not set")
    else:
        print(f"budget_mode: {budget}")
    print("next artifacts:")
    for artifact in artifacts.values():
        print(f"- {artifact}")
    return 0


SLIDE_ID_PATTERN = re.compile(r"^\d{2}_[a-z0-9_]+$")
SLIDE_HEADING_PATTERN = re.compile(r"^##\s+(\S+)\s*$", re.MULTILINE)

TITLE_LENGTH_MAX = 80

KNOWN_ARCHETYPES = {
    "cover", "thesis", "timeline", "comparison", "evidence cards",
    "process", "map", "data chart", "tension", "closing",
}

SUPPORT_COUNT_BANDS: dict[str, tuple[int, int] | None] = {
    "cover": None,
    "closing": None,
    "process": (2, 7),
    "evidence cards": (2, 10),
    "timeline": (2, 10),
    "thesis": (2, 6),
}
SUPPORT_COUNT_DEFAULT = (2, 4)

SLIDE_COUNT_BANDS = {
    "quick": (5, 7),
    "balanced": (7, 10),
    "premium": (8, 14),
}

NATIVE_PPTX_SCRIPT_PATTERNS: tuple[tuple[str, str], ...] = (
    (r"^\s*(from|import)\s+pptx\b", "imports python-pptx"),
    (r"\bPresentation\s*\(", "constructs a PowerPoint presentation"),
    (r"\bslide\.shapes\b", "uses native PowerPoint slide shapes"),
    (r"\bshapes\.add_(textbox|chart|table|shape|connector)\b", "adds native PowerPoint objects"),
    (r"\btext_frame\b", "edits native PowerPoint text frames"),
    (r"\bMSO_SHAPE\b|\bXL_CHART_TYPE\b", "uses native PowerPoint shape/chart constants"),
)

FORBIDDEN_PROMPT_PATTERNS: tuple[tuple[str, str], ...] = (
    (r"\bnative[- ]pptx\b", "asks for native-PPTX output"),
    (r"\bpptx[- ]native\b", "asks for PPTX-native output"),
    (r"\beditable\s+(powerpoint|pptx|ppt)\b", "asks for editable PowerPoint output"),
    (r"\b(powerpoint|pptx|ppt)\s+(text boxes|shapes|charts|layouts)\b", "asks for native PowerPoint objects"),
    (r"\bshape[- ]by[- ]shape\b", "asks for shape-by-shape slide construction"),
    (r"\bpython[- ]pptx\b", "mentions python-pptx in a production artifact"),
    (r"\bhybrid\s+mode\b", "asks for disabled hybrid mode"),
    (r"\buse\s+(codex\s+)?presentations\b", "delegates to a native presentation tool"),
)

NATIVE_PPTX_NEGATION_PATTERN = re.compile(
    r"\b(no|not|never|without|forbidden|disabled|unsupported|does not|do not|is not|must not)\b|不支持|不要|不能|禁止|不是",
    re.IGNORECASE,
)


def _scan_brief(brief_path: Path) -> list[tuple[str, str, str]]:
    findings: list[tuple[str, str, str]] = []
    if not brief_path.is_file():
        findings.append(("error", "BRIEF-001", "work/deck-brief.md missing"))
        return findings
    text = brief_path.read_text(encoding="utf-8").lower()
    for token, rule_id, label in (
        ("thesis", "BRIEF-THESIS", "thesis"),
        ("audience", "BRIEF-AUDIENCE", "audience"),
        ("arc", "BRIEF-ARC", "arc"),
    ):
        if not re.search(rf"^##\s+[^\n]*{token}", text, re.MULTILINE):
            findings.append(("warn", rule_id, f"deck-brief.md missing a heading mentioning '{label}'"))
    return findings


def _parse_storyboard(sb_path: Path) -> tuple[list[tuple[str, str]], list[tuple[str, str, str]]]:
    """Return (slide_blocks, findings). Each slide_block is (slide_id, body_text)."""
    findings: list[tuple[str, str, str]] = []
    if not sb_path.is_file():
        findings.append(("error", "SB-001", "work/storyboard.md missing"))
        return [], findings
    text = sb_path.read_text(encoding="utf-8")
    slide_blocks: list[tuple[str, str]] = []
    matches = list(SLIDE_HEADING_PATTERN.finditer(text))
    if not matches:
        findings.append(("error", "SB-002", "no slide sections (## headings) found in storyboard"))
        return [], findings
    for idx, match in enumerate(matches):
        slide_id = match.group(1)
        if slide_id.lower() in {"deck", "meta"} or slide_id.startswith("Deck"):
            continue
        body_start = match.end()
        body_end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        body = text[body_start:body_end]
        if not SLIDE_ID_PATTERN.match(slide_id):
            findings.append(("warn", "SLIDE-ID-FORMAT", f"slide id '{slide_id}' does not match NN_slug format"))
        slide_blocks.append((slide_id, body))
    return slide_blocks, findings


def _extract_field_value(body: str, name: str) -> str:
    pattern = rf"^-\s+\*\*{re.escape(name)}\*\*:\s*(.+)$"
    m = re.search(pattern, body, re.MULTILINE | re.IGNORECASE)
    return m.group(1).strip() if m else ""


def _normalize_archetype(value: str) -> str:
    """Strip parenthetical clarifications and lowercase. 'Thesis (four pillars)' -> 'thesis'."""
    cleaned = re.sub(r"\([^)]*\)", "", value).strip().lower()
    return cleaned


def _scan_slide(slide_id: str, body: str) -> list[tuple[str, str, str]]:
    findings: list[tuple[str, str, str]] = []
    body_lower = body.lower()
    for token, rule_id, label in (
        ("title", "SLIDE-TITLE", "Title"),
        ("primary job", "SLIDE-JOB", "Primary job"),
        ("core claim", "SLIDE-CLAIM", "Core claim"),
        ("support", "SLIDE-SUPPORT", "Support points"),
    ):
        if token not in body_lower:
            findings.append(("warn", rule_id, f"slide {slide_id}: missing '{label}' field"))

    title_value = _extract_field_value(body, "Title")
    if title_value and len(title_value) > TITLE_LENGTH_MAX:
        findings.append(
            (
                "warn",
                "SLIDE-TITLE-LENGTH",
                f"slide {slide_id}: title length {len(title_value)} exceeds {TITLE_LENGTH_MAX} chars",
            )
        )

    archetype_value = _extract_field_value(body, "Archetype")
    archetype_norm = _normalize_archetype(archetype_value)
    if archetype_value and archetype_norm not in KNOWN_ARCHETYPES:
        findings.append(
            (
                "warn",
                "SLIDE-ARCHETYPE-UNKNOWN",
                f"slide {slide_id}: archetype '{archetype_value}' is not in slide-archetypes.md",
            )
        )

    sp_match = re.search(
        r"\*\*support points?\*\*[^\n]*\n((?:[ \t]*[-*][ \t]+[^\n]*\n)+)",
        body,
        re.IGNORECASE,
    )
    if sp_match:
        bullets = re.findall(r"^[ \t]*[-*][ \t]+", sp_match.group(1), re.MULTILINE)
        band_key = archetype_norm if archetype_norm in SUPPORT_COUNT_BANDS else None
        band = SUPPORT_COUNT_BANDS.get(band_key, SUPPORT_COUNT_DEFAULT) if band_key is not None else SUPPORT_COUNT_DEFAULT
        if band is None:
            return findings
        low, high = band
        if not (low <= len(bullets) <= high):
            findings.append(
                (
                    "warn",
                    "SLIDE-SUPPORT-COUNT",
                    f"slide {slide_id} ({archetype_norm or 'no-archetype'}): support points count {len(bullets)} outside {low}-{high} range",
                )
            )
    return findings


def _scan_slide_count(budget: str | None, slide_count: int) -> list[tuple[str, str, str]]:
    if not budget:
        return []
    band = SLIDE_COUNT_BANDS.get(budget)
    if band is None:
        return []
    low, high = band
    if low <= slide_count <= high:
        return []
    return [
        (
            "warn",
            "DECK-SLIDE-COUNT",
            f"slide count {slide_count} is outside the {budget} band ({low}-{high})",
        )
    ]


def _scan_prompts(run_dir: Path, storyboard_ids: list[str]) -> list[tuple[str, str, str]]:
    findings: list[tuple[str, str, str]] = []
    prompts_dir = run_dir / "prompts"
    if not prompts_dir.is_dir():
        findings.append(("error", "PROMPTS-001", "prompts/ directory missing"))
        return findings
    prompt_ids = sorted(p.stem for p in prompts_dir.glob("*.md") if p.name != "README.md")
    storyboard_set = set(storyboard_ids)
    for sid in storyboard_ids:
        if sid not in prompt_ids:
            findings.append(("error", "PROMPT-MISSING", f"slide {sid} has no prompt file at prompts/{sid}.md"))
    for pid in prompt_ids:
        if pid not in storyboard_set:
            findings.append(("warn", "PROMPT-ORPHAN", f"prompt prompts/{pid}.md has no matching storyboard slide"))
    return findings


def _scan_native_pptx_language(run_dir: Path) -> list[tuple[str, str, str]]:
    """Reject production artifacts that steer an image-first run back to native-PPTX."""
    findings: list[tuple[str, str, str]] = []
    candidate_files: list[Path] = []

    storyboard_path = run_dir / "work" / "storyboard.md"
    if storyboard_path.is_file():
        candidate_files.append(storyboard_path)

    prompts_dir = run_dir / "prompts"
    if prompts_dir.is_dir():
        candidate_files.extend(sorted(prompts_dir.glob("*.md")))

    for path in candidate_files:
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            continue
        for line_number, line in enumerate(text.splitlines(), start=1):
            if NATIVE_PPTX_NEGATION_PATTERN.search(line):
                continue
            for pattern, reason in FORBIDDEN_PROMPT_PATTERNS:
                if not re.search(pattern, line, re.IGNORECASE):
                    continue
                findings.append(
                    (
                        "error",
                        "IMG-FIRST-NATIVE-PPTX-LANGUAGE",
                        f"{path.relative_to(run_dir)}:{line_number} {reason}; Deckit prompts/storyboards must stay image-first",
                    )
                )
                break
    return findings


def _scan_native_pptx_scripts(run_dir: Path) -> list[tuple[str, str, str]]:
    """Reject run-local scripts that look like native PowerPoint assembly."""
    findings: list[tuple[str, str, str]] = []
    scripts_dir = run_dir / "scripts"
    if not scripts_dir.is_dir():
        return findings

    script_suffixes = {".py", ".js", ".ts", ".mjs", ".cjs"}
    for script_path in sorted(p for p in scripts_dir.rglob("*") if p.is_file() and p.suffix.lower() in script_suffixes):
        try:
            text = script_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            continue

        matched_reasons = [
            reason
            for pattern, reason in NATIVE_PPTX_SCRIPT_PATTERNS
            if re.search(pattern, text, re.IGNORECASE | re.MULTILINE)
        ]
        if matched_reasons:
            findings.append(
                (
                    "error",
                    "IMG-FIRST-NATIVE-PPTX-SCRIPT",
                    f"{script_path.relative_to(run_dir)} {', '.join(matched_reasons)}; PPTX is only a container after $imagegen PNGs exist, never a native production route",
                )
            )
    return findings


def _scan_image_first_artifacts(run_dir: Path, storyboard_ids: list[str]) -> list[tuple[str, str, str]]:
    findings: list[tuple[str, str, str]] = []
    generated_dir = run_dir / "assets" / "generated-slides"
    generated_pngs = sorted(generated_dir.glob("*.png")) if generated_dir.is_dir() else []
    generated_ids = {p.stem for p in generated_pngs}
    dist_pptx = sorted((run_dir / "dist").glob("*.pptx")) if (run_dir / "dist").is_dir() else []

    if dist_pptx and not generated_pngs:
        pptx_names = ", ".join(p.name for p in dist_pptx)
        findings.append(
            (
                "error",
                "IMG-FIRST-PPTX-WITHOUT-GENERATED-SLIDES",
                f"dist contains PPTX deliverable(s) ({pptx_names}) but assets/generated-slides has no PNGs from $imagegen",
            )
        )
    elif dist_pptx and storyboard_ids and set(storyboard_ids) - generated_ids:
        missing = ", ".join(sorted(set(storyboard_ids) - generated_ids))
        pptx_names = ", ".join(p.name for p in dist_pptx)
        findings.append(
            (
                "error",
                "IMG-FIRST-PPTX-BEFORE-COMPLETE-IMAGES",
                f"dist contains PPTX deliverable(s) ({pptx_names}) before every storyboard slide has a generated PNG; missing: {missing}",
            )
        )

    for sid in storyboard_ids:
        if generated_pngs and sid not in generated_ids:
            findings.append(
                (
                    "error",
                    "IMG-FIRST-GENERATED-SLIDE-MISSING",
                    f"slide {sid} has no generated image at assets/generated-slides/{sid}.png",
                )
            )
    for image_id in sorted(generated_ids - set(storyboard_ids)):
        findings.append(
            (
                "warn",
                "IMG-FIRST-GENERATED-SLIDE-ORPHAN",
                f"generated image assets/generated-slides/{image_id}.png has no matching storyboard slide",
            )
        )

    findings.extend(_scan_native_pptx_scripts(run_dir))
    findings.extend(_scan_native_pptx_language(run_dir))
    return findings


def _package_images(run_dir: Path, out_path: Path | None) -> int:
    """Package generated slide PNGs into a non-editable image-only PPTX container."""
    run_dir = run_dir.resolve()
    if not run_dir.is_dir():
        raise NotADirectoryError(f"run directory does not exist: {run_dir}")

    slide_blocks, sb_findings = _parse_storyboard(run_dir / "work" / "storyboard.md")
    errors = [finding for finding in sb_findings if finding[0] == "error"]
    if errors:
        detail = "; ".join(message for _, _, message in errors)
        raise ValueError(f"cannot package images because storyboard is invalid: {detail}")

    slide_ids = [slide_id for slide_id, _ in slide_blocks]
    if not slide_ids:
        raise ValueError("cannot package images because storyboard contains no slides")

    generated_dir = run_dir / "assets" / "generated-slides"
    missing = [slide_id for slide_id in slide_ids if not (generated_dir / f"{slide_id}.png").is_file()]
    if missing:
        raise FileNotFoundError(
            "cannot package PPTX before every storyboard slide has a generated PNG. Missing: "
            + ", ".join(f"assets/generated-slides/{slide_id}.png" for slide_id in missing)
        )

    if out_path is None:
        out_path = run_dir / "dist" / f"{run_dir.name}.pptx"
    elif not out_path.is_absolute():
        out_path = (run_dir / out_path).resolve()
    else:
        out_path = out_path.resolve()

    dist_dir = run_dir / "dist"
    dist_dir.mkdir(parents=True, exist_ok=True)
    _ensure_inside(dist_dir, out_path)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for stable PPTX container packaging. "
            "Install the dev tools with their dependencies, then rerun deckit-dev package-images."
        ) from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for slide_id in slide_ids:
        slide = prs.slides.add_slide(blank_layout)
        image_path = generated_dir / f"{slide_id}.png"
        # Packaging discipline: one generated full-slide image per slide; no editable text boxes/charts/native layouts.
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(str(out_path))

    # Structural smoke checks: the file is a readable ZIP, re-opens through python-pptx, and has one picture per slide.
    with zipfile.ZipFile(out_path) as package:
        bad_member = package.testzip()
    if bad_member is not None:
        raise ValueError(f"packaged PPTX failed ZIP integrity check at {bad_member}")

    check = Presentation(str(out_path))
    if len(check.slides) != len(slide_ids):
        raise ValueError(f"packaged PPTX has {len(check.slides)} slide(s), expected {len(slide_ids)}")
    for idx, slide in enumerate(check.slides, start=1):
        if len(slide.shapes) != 1:
            raise ValueError(f"packaged PPTX slide {idx} has {len(slide.shapes)} shape(s), expected exactly 1 image")

    print(f"wrote: {out_path}")
    print(f"slides packaged: {len(slide_ids)}")
    print("mode: image-first PNG container; no editable native PowerPoint content")
    return 0


def review(run_dir: Path) -> int:
    run_dir = run_dir.resolve()
    if not run_dir.is_dir():
        raise NotADirectoryError(f"run directory does not exist: {run_dir}")

    findings: list[tuple[str, str, str]] = []

    run_json_path = run_dir / "run.json"
    manifest: dict[str, object] = {}
    if not run_json_path.is_file():
        findings.append(("error", "RUN-001", "run.json missing"))
    else:
        try:
            manifest = json.loads(run_json_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            findings.append(("error", "RUN-002", f"run.json is not valid JSON: {exc}"))

    mode = manifest.get("production_mode") if isinstance(manifest, dict) else None
    if mode is None:
        findings.append(("warn", "MODE-001", "production_mode is null in run.json"))
    elif mode not in PRODUCTION_MODES:
        findings.append(("error", "MODE-002", f"production_mode '{mode}' is not one of {PRODUCTION_MODES}"))

    findings.extend(_scan_brief(run_dir / "work" / "deck-brief.md"))

    slide_blocks, sb_findings = _parse_storyboard(run_dir / "work" / "storyboard.md")
    findings.extend(sb_findings)
    for slide_id, body in slide_blocks:
        findings.extend(_scan_slide(slide_id, body))

    storyboard_ids = [sid for sid, _ in slide_blocks]
    budget = manifest.get("budget_mode") if isinstance(manifest, dict) else None
    if storyboard_ids and isinstance(budget, str):
        findings.extend(_scan_slide_count(budget, len(storyboard_ids)))
    findings.extend(_scan_prompts(run_dir, storyboard_ids))
    if mode == "image-first":
        findings.extend(_scan_image_first_artifacts(run_dir, storyboard_ids))

    counts = {"error": 0, "warn": 0}
    for sev, _, _ in findings:
        counts[sev] = counts.get(sev, 0) + 1

    dist_dir = run_dir / "dist"
    dist_dir.mkdir(parents=True, exist_ok=True)
    review_path = dist_dir / "review.md"

    lines: list[str] = []
    lines.append(f"# Review Report — {run_dir.name}\n\n")
    lines.append(f"- Production mode: `{mode if mode else 'unset'}`\n")
    lines.append(f"- Slides found in storyboard: {len(storyboard_ids)}\n")
    lines.append(f"- Findings: {counts['error']} error(s), {counts['warn']} warning(s)\n\n")
    if not findings:
        lines.append("All checks passed.\n")
    else:
        lines.append("## Findings\n\n")
        for sev, rule_id, msg in findings:
            lines.append(f"- **{sev.upper()}** `{rule_id}` — {msg}\n")
    review_path.write_text("".join(lines), encoding="utf-8")

    print(f"wrote: {review_path}")
    print(f"errors: {counts['error']}; warnings: {counts['warn']}")
    return 0 if counts["error"] == 0 else 1


def main() -> int:
    parser = argparse.ArgumentParser(description="Deckit development tools")
    subparsers = parser.add_subparsers(dest="command", required=True)

    inspect_parser = subparsers.add_parser("inspect", help="Inspect the plugin manifest and skills")
    inspect_parser.add_argument("--plugin", type=Path, default=DEFAULT_PLUGIN)

    marketplace_parser = subparsers.add_parser("inspect-marketplace", help="Inspect repo-local plugin marketplace")
    marketplace_parser.add_argument("--marketplace", type=Path, default=DEFAULT_MARKETPLACE)

    new_run_parser = subparsers.add_parser("new-run", help="Create a standard local run directory from a text, PDF, or URL source")
    new_run_parser.add_argument("--source", required=True, help="Path to a text/Markdown file, a .pdf file, or an http(s) URL")
    new_run_parser.add_argument("--name", help="Run name. Defaults to the source filename stem (or URL host+path).")
    new_run_parser.add_argument("--runs-dir", type=Path, default=DEFAULT_RUNS_DIR)
    new_run_parser.add_argument("--force", action="store_true", help="Reuse an existing run directory and replace source copy")
    new_run_parser.add_argument(
        "--mode",
        choices=PRODUCTION_MODES,
        default="image-first",
        help="Deprecated compatibility flag. Deckit always uses image-first; no other production modes are valid.",
    )
    new_run_parser.add_argument("--budget", choices=("quick", "balanced", "premium"), help="Budget mode (recorded in run.json).")

    review_parser = subparsers.add_parser("review", help="Run rule-based quality checks on a run folder and write dist/review.md")
    review_parser.add_argument("--run", type=Path, required=True, help="Path to the run directory to review.")

    package_parser = subparsers.add_parser(
        "package-images",
        help="Package assets/generated-slides/*.png into a non-editable image-only PPTX container",
    )
    package_parser.add_argument("--run", type=Path, required=True, help="Path to the run directory to package.")
    package_parser.add_argument(
        "--out",
        type=Path,
        help="Output PPTX path. Relative paths are resolved inside the run directory; default: dist/<run-name>.pptx.",
    )

    ingest_parser = subparsers.add_parser("ingest", help="Convert a PDF or URL into source/input.md (Markdown)")
    ingest_group = ingest_parser.add_mutually_exclusive_group(required=True)
    ingest_group.add_argument("--pdf", type=Path, help="Path to a .pdf file")
    ingest_group.add_argument("--url", help="An http(s) URL")
    ingest_parser.add_argument("--out", type=Path, required=True, help="Output Markdown path (e.g. source/input.md)")

    args = parser.parse_args()

    if args.command == "inspect":
        return inspect_plugin(args.plugin.resolve())
    if args.command == "inspect-marketplace":
        return inspect_marketplace(args.marketplace)
    if args.command == "new-run":
        return new_run(args.source, args.name, args.runs_dir, args.force, args.mode, args.budget)
    if args.command == "review":
        return review(args.run)
    if args.command == "package-images":
        return _package_images(args.run, args.out)
    if args.command == "ingest":
        from deckit_dev.ingest import ingest_pdf, ingest_url

        if args.pdf is not None:
            return ingest_pdf(args.pdf, args.out)
        return ingest_url(args.url, args.out)

    parser.error(f"unknown command: {args.command}")
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
